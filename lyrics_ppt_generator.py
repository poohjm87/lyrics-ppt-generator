#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import os
import re
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import pptx


def read_lyrics_files(source_dir):
    """source 디렉토리에서 txt 파일들을 읽어 가사 데이터를 반환"""
    lyrics_data = []
    
    if not os.path.exists(source_dir):
        print(f"디렉토리가 존재하지 않습니다: {source_dir}")
        return lyrics_data
    
    # txt 파일들을 순서대로 정렬
    txt_files = [f for f in os.listdir(source_dir) if f.endswith('.txt')]
    txt_files.sort()
    
    for filename in txt_files:
        filepath = os.path.join(source_dir, filename)
        
        # 파일명에서 순서와 제목 추출
        match = re.match(r'(\d+)\.(.+)\.txt', filename)
        if match:
            order = match.group(1)
            title = match.group(2)
        else:
            title = filename.replace('.txt', '')
        
        try:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as file:
                content = file.read().strip()
                lyrics_lines = [line.strip() for line in content.split('\n') if line.strip()]
                
                lyrics_data.append({
                    'title': title,
                    'lyrics': lyrics_lines
                })
                print(f"읽은 파일: {filename} ({len(lyrics_lines)}줄)")
        
        except UnicodeDecodeError as e:
            print(f"파일 인코딩 오류 {filename}: UTF-8로 저장해주세요")
        except Exception as e:
            print(f"파일 읽기 오류 {filename}: {e}")
    
    return lyrics_data


def create_lyrics_ppt(lyrics_data, output_filename, template_path=None):
    """가사 데이터를 이용해 PPT 파일 생성"""
    try:
        # 템플릿 파일이 있으면 사용, 없거나 오류가 발생하면 기본 프레젠테이션 생성
        if template_path and os.path.exists(template_path):
            try:
                print(f"템플릿 파일 사용: {template_path}")
                prs = Presentation(template_path)
            except Exception as e:
                print(f"템플릿 파일 로드 실패: {e}")
                print("기본 템플릿 사용")
                prs = Presentation()
        else:
            print("기본 템플릿 사용")
            prs = Presentation()
    except Exception as e:
        print(f"프레젠테이션 생성 실패: {e}")
        return False
    
    try:
        # 제목 슬라이드용 레이아웃 (Title Slide)
        title_layout = prs.slide_layouts[0]
        # 가사용 레이아웃 (Title and Content)
        content_layout = prs.slide_layouts[1]
        # 빈 슬라이드용 레이아웃
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
        
        for song_idx, song in enumerate(lyrics_data):
            print(f"처리 중: {song['title']}")
            
            try:
                # 가사를 2줄씩 슬라이드에 추가 (Title and Content 레이아웃 사용)
                lyrics = song['lyrics']
                for i in range(0, len(lyrics), 2):
                    content_slide = prs.slides.add_slide(content_layout)
                    
                    # 제목 영역에 노래 제목 추가
                    content_slide.shapes.title.text = song['title']
                    
                    # 내용 영역에 가사 추가
                    content_shape = content_slide.shapes.placeholders[1]
                    text_frame = content_shape.text_frame
                    text_frame.clear()  # 기본 텍스트 제거
                    
                    # 첫 번째 줄
                    p1 = text_frame.paragraphs[0]
                    p1.text = lyrics[i]
                    p1.alignment = PP_ALIGN.CENTER
                    
                    # 두 번째 줄 (있는 경우)
                    if i + 1 < len(lyrics):
                        p2 = text_frame.add_paragraph()
                        p2.text = lyrics[i + 1]
                        p2.alignment = PP_ALIGN.CENTER
                
                # 노래와 노래 사이에 빈 페이지 추가 (마지막 노래 제외)
                if song_idx < len(lyrics_data) - 1:
                    blank_slide = prs.slides.add_slide(blank_layout)
                    
            except Exception as e:
                print(f"슬라이드 생성 중 오류 ({song['title']}): {e}")
                continue
        
        # PPT 파일 저장
        try:
            print(f"PPT 파일 저장 중: {output_filename}")
            print(f"현재 작업 디렉토리: {os.getcwd()}")
            print(f"출력 파일 전체 경로: {os.path.abspath(output_filename)}")
            
            # 파일이 이미 존재하면 삭제
            if os.path.exists(output_filename):
                print(f"기존 파일 삭제: {output_filename}")
                os.remove(output_filename)
            
            prs.save(output_filename)
            
            # 저장 확인
            if os.path.exists(output_filename):
                print(f"PPT 파일이 생성되었습니다: {output_filename}")
                return True
            else:
                print("파일이 저장되지 않았습니다.")
                return False
                
        except PermissionError as e:
            print(f"권한 오류: {e}")
            print("다른 프로그램에서 파일이 열려있거나 관리자 권한이 필요할 수 있습니다.")
            return False
        except Exception as e:
            print(f"PPT 파일 저장 실패: {e}")
            print("다음을 확인해주세요:")
            print("1. 디스크 용량이 충분한지")
            print("2. 파일명에 특수문자가 없는지") 
            print("3. 폴더에 쓰기 권한이 있는지")
            return False
            
    except Exception as e:
        print(f"PPT 생성 중 전체 오류: {e}")
        return False


def main():
    """메인 함수"""
    print("=== 가사 PPT 생성기 디버깅 정보 ===")
    print(f"Python 버전: {sys.version}")
    print(f"python-pptx 버전: {pptx.__version__}")
    print(f"현재 작업 디렉토리: {os.getcwd()}")
    print(f"스크립트 위치: {os.path.abspath(__file__)}")
    
    # 경로 설정
    current_dir = os.getcwd()
    source_dir = os.path.join(current_dir, "source")
    
    # 바탕화면에 저장 (확실한 쓰기 권한)
    desktop_dir = os.path.join(os.path.expanduser("~"), "Desktop")
    
    # 출력 디렉토리 결정 (우선순위: 현재 디렉토리 -> 바탕화면 -> 임시폴더)
    output_dir = current_dir
    
    # 현재 디렉토리에 쓰기 권한 테스트
    try:
        test_file = os.path.join(current_dir, "write_test.tmp")
        with open(test_file, 'w') as f:
            f.write("test")
        os.remove(test_file)
        print("✓ 현재 디렉토리에 저장됩니다")
    except:
        # 바탕화면 사용
        if os.path.exists(desktop_dir):
            output_dir = desktop_dir
            print(f"✓ 바탕화면에 저장됩니다: {output_dir}")
        else:
            # 임시 폴더 사용
            import tempfile
            output_dir = tempfile.gettempdir()
            print(f"✓ 임시 폴더에 저장됩니다: {output_dir}")
    
    # 현재 날짜로 파일명 생성 (Windows 호환 경로)
    today = datetime.now()
    filename = f"{today.strftime('%Y%m%d')}_Lyrics.pptx"
    output_filename = os.path.normpath(os.path.join(output_dir, filename))
    
    # 템플릿 파일 경로 확인
    template_path = os.path.join(source_dir, "template.pptx")
    
    print("=== 가사 PPT 생성기 ===")
    print(f"소스 디렉토리: {source_dir}")
    print(f"출력 파일: {output_filename}")
    
    # 템플릿 파일 존재 여부 확인
    if os.path.exists(template_path):
        print(f"템플릿 파일 발견: {template_path}")
    else:
        print("템플릿 파일이 없습니다. 기본 템플릿을 사용합니다.")
        template_path = None  # 파일이 없으면 None으로 설정
    
    # 가사 파일 읽기
    lyrics_data = read_lyrics_files(source_dir)
    
    if not lyrics_data:
        print("가사 파일을 찾을 수 없습니다.")
        return
    
    print(f"총 {len(lyrics_data)}개의 노래를 찾았습니다.")
    
    # PPT 생성 (템플릿 파일 경로 전달)
    success = create_lyrics_ppt(lyrics_data, output_filename, template_path)
    
    if success:
        print("작업이 완료되었습니다!")
        print(f"생성된 파일: {os.path.abspath(output_filename)}")
    else:
        print("PPT 파일 생성에 실패했습니다.")
        print("다음 사항을 확인해주세요:")
        print("1. 현재 폴더에 쓰기 권한이 있는지")
        print("2. 동일한 이름의 파일이 다른 프로그램에서 열려있지 않은지")
        print("3. 디스크 용량이 충분한지")
    
    # Windows에서 창이 바로 닫히지 않도록 대기
    if sys.platform.startswith('win'):
        input("엔터 키를 누르면 종료됩니다...")


if __name__ == "__main__":
    main()